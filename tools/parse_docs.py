#!/usr/bin/env python3
"""
parse_docs.py — Extract text from project folders and/or public URLs.

Usage:
    # 1. Folders only (each folder = one project, id = slug of folder name)
    python tools/parse_docs.py <folder1> [<folder2> ...] --output work/parsed.json

    # 2. Mixed folders + URLs via a JSON spec file
    python tools/parse_docs.py --spec projects.spec.json --output work/parsed.json

Spec file format (see examples/projects.spec.example.json):
    {
      "projects": [
        {
          "id": "acme-support",                  # kebab_case slug
          "name": "Acme Support Assistant",       # display name
          "inputs": [                             # folder paths OR public URLs
            "./projects/acme-support",
            "https://arxiv.org/abs/2606.12345"
          ]
        }
      ]
    }

Authenticated URLs (SharePoint / Confluence / Notion / 飞书 / Google Docs …)
are NOT supported directly — HTTP 401/403 is reported as status='auth_required'.
When the agent sees this, it should ask the user to either:
  (a) download the page/file in a browser and put it in a local folder input, or
  (b) use an MCP connector (SharePoint / Confluence / Notion MCP) that has
      authenticated access to fetch the content, saving to work/downloads/ then
      referencing that folder as an input.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import Request, urlopen
from urllib.error import HTTPError, URLError

SUPPORTED_EXT = {".pptx", ".pdf", ".xlsx", ".xls", ".docx", ".md", ".txt"}
DEFAULT_UA = (
    "Mozilla/5.0 (compatible; ai-brief/1.0) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/120.0 Safari/537.36"
)
AUTH_ADVICE = (
    "URL 需要认证。请让用户 (a) 在浏览器里下载后放到本地文件夹作为 folder 输入，"
    "或 (b) 用支持该服务的 MCP（SharePoint/Confluence/Notion 等）抓下来存到 "
    "work/downloads/ 目录作为 folder 输入。"
)


def slugify(name: str) -> str:
    s = re.sub(r"[^\w一-鿿\-]+", "_", name, flags=re.UNICODE).strip("_").lower()
    return s or "project"


# ---------------- File loaders ----------------

def _pptx(f: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(f))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"=== Slide {i} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip()
                    if txt:
                        parts.append(txt)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_txt = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_txt:
                        parts.append(row_txt)
        chunks.append("\n".join(parts))
    return "\n\n".join(chunks)


def _pdf(f: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(f))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _xlsx(f: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(f), read_only=True, data_only=True)
    out = []
    for name in wb.sheetnames:
        ws = wb[name]
        out.append(f"=== Sheet: {name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                out.append(" | ".join(cells))
    return "\n".join(out)


def _docx(f: Path) -> str:
    from docx import Document
    doc = Document(str(f))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_txt = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_txt:
                parts.append(row_txt)
    return "\n".join(parts)


LOADERS = {".pptx": _pptx, ".pdf": _pdf, ".xlsx": _xlsx, ".xls": _xlsx, ".docx": _docx}


def _extract_local(f: Path) -> str:
    ext = f.suffix.lower()
    if ext in LOADERS:
        return LOADERS[ext](f)
    return f.read_text(encoding="utf-8", errors="replace")


# ---------------- URL fetch & HTML extraction ----------------

def _fetch_url(url: str, timeout: float = 20.0) -> dict:
    """Return {status, http?, content_type?, body?, error?, advice?}.

    status: 'ok' | 'auth_required' | 'not_found' | 'error'
    """
    headers = {"User-Agent": DEFAULT_UA, "Accept": "text/html,application/xhtml+xml,*/*"}
    req = Request(url, headers=headers)
    try:
        with urlopen(req, timeout=timeout) as resp:
            ct = resp.headers.get("Content-Type", "").lower()
            body = resp.read()
            return {"status": "ok", "http": resp.status, "content_type": ct, "body": body}
    except HTTPError as e:
        if e.code in (401, 403):
            return {"status": "auth_required", "http": e.code, "advice": AUTH_ADVICE}
        if e.code == 404:
            return {"status": "not_found", "http": 404, "advice": "URL 返回 404，请核对链接是否失效"}
        return {"status": "error", "http": e.code, "error": str(e)}
    except (URLError, TimeoutError) as e:
        return {"status": "error", "error": str(e), "advice": "网络错误：请稍后重试或检查代理"}
    except Exception as e:
        return {"status": "error", "error": str(e)}


def _extract_html_text(body: bytes) -> str:
    """Extract readable text from HTML using BeautifulSoup (already in requirements).

    Falls back to raw text if BS4 unavailable.
    """
    try:
        text = body.decode("utf-8", errors="replace")
    except Exception:
        text = str(body)

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return re.sub(r"<[^>]+>", " ", text)

    soup = BeautifulSoup(text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "aside", "noscript"]):
        tag.decompose()

    main = soup.find("article") or soup.find("main") or soup.body or soup
    lines = []
    title = soup.find("title")
    if title and title.string:
        lines.append(f"# {title.string.strip()}\n")
    for el in main.find_all(["h1", "h2", "h3", "h4", "p", "li", "pre", "code", "blockquote"]):
        s = el.get_text(" ", strip=True)
        if not s:
            continue
        tag = el.name
        if tag.startswith("h"):
            lines.append(f"\n## {s}\n")
        elif tag == "li":
            lines.append(f"- {s}")
        else:
            lines.append(s)
    result = "\n".join(lines).strip()
    if not result:
        result = main.get_text("\n", strip=True)
    return result


# ---------------- Project assembly ----------------

def _parse_folder_source(folder: Path) -> tuple[dict, list[dict]]:
    """Return (source_record, files[])."""
    if not folder.exists():
        return ({"kind": "folder", "path": str(folder), "status": "not_found",
                 "advice": f"路径不存在：{folder}"}, [])
    if not folder.is_dir():
        return ({"kind": "folder", "path": str(folder), "status": "error",
                 "advice": f"不是目录：{folder}"}, [])
    files_out = []
    for f in sorted(folder.rglob("*")):
        if not f.is_file():
            continue
        if f.suffix.lower() not in SUPPORTED_EXT:
            continue
        try:
            text = _extract_local(f)
        except Exception as e:
            print(f"[parse_docs] {f.name}: {e}", file=sys.stderr)
            continue
        if not text or not text.strip():
            continue
        files_out.append({
            "name": str(f.relative_to(folder)),
            "type": f.suffix.lower().lstrip("."),
            "source": "folder",
            "path": str(f.relative_to(folder)),
            "text": text,
        })
    return ({"kind": "folder", "path": str(folder.resolve()), "status": "ok",
             "file_count": len(files_out)}, files_out)


def _parse_url_source(url: str) -> tuple[dict, list[dict]]:
    """Return (source_record, files[])."""
    result = _fetch_url(url)
    record = {"kind": "url", "url": url, "status": result["status"]}
    if "http" in result: record["http"] = result["http"]
    if "content_type" in result: record["content_type"] = result["content_type"]
    if "advice" in result: record["advice"] = result["advice"]
    if "error" in result: record["error"] = result["error"]

    if result["status"] != "ok":
        return record, []

    body = result["body"]
    ct = result.get("content_type", "")
    if "html" in ct or "xml" in ct or not ct:
        text = _extract_html_text(body)
    else:
        try:
            text = body.decode("utf-8", errors="replace")
        except Exception:
            text = str(body)

    if not text.strip():
        record["status"] = "empty"
        record["advice"] = "URL 可达但抽不到正文（可能是 JS 渲染页或空页）。请下载到本地或换一个直接的 URL。"
        return record, []

    name = _url_display_name(url)
    return record, [{
        "name": name,
        "type": "url",
        "source": "url",
        "url": url,
        "text": text,
    }]


def _url_display_name(url: str) -> str:
    p = urlparse(url)
    path_part = re.sub(r"[^\w\-\.]+", "_", p.path.strip("/"))[:60] or "root"
    return f"{p.hostname}_{path_part}"


def parse_project(spec: dict) -> dict:
    """spec = {'id': ..., 'name': ..., 'inputs': [str, ...]}

    Each input is a string — either a filesystem path or an http(s):// URL.
    """
    proj = {"id": spec["id"], "name": spec.get("name") or spec["id"],
            "sources": [], "files": []}
    for inp in spec.get("inputs", []):
        if _is_url(inp):
            src, files = _parse_url_source(inp)
        else:
            src, files = _parse_folder_source(Path(inp).expanduser())

        proj["sources"].append(src)
        proj["files"].extend(files)
        if src["status"] == "ok":
            n = len(files) if src["kind"] == "url" else f"{src.get('file_count', 0)} file(s)"
            print(f"[parse_docs] ✓ [{spec['id']}] {src['kind']}: {inp} → {n}",
                  file=sys.stderr)
        else:
            print(f"[parse_docs] ✗ [{spec['id']}] {src['kind']}: {inp} → "
                  f"{src['status']}: {src.get('advice', '')}", file=sys.stderr)
    return proj


def _is_url(s: str) -> bool:
    return isinstance(s, str) and (s.startswith("http://") or s.startswith("https://"))


# ---------------- CLI ----------------

def _spec_from_positional(paths: list[Path]) -> list[dict]:
    """Backward-compat: each folder becomes its own project."""
    out = []
    for p in paths:
        if _is_url(str(p)):
            print(f"[parse_docs] URL positional arg without --spec: {p}. "
                  f"Wrap it in a spec file.", file=sys.stderr)
            sys.exit(2)
        p = Path(p).expanduser()
        out.append({"id": slugify(p.name), "name": p.name, "inputs": [str(p)]})
    return out


def _load_spec(spec_path: Path) -> list[dict]:
    data = json.loads(spec_path.read_text(encoding="utf-8"))
    projects = data.get("projects") or []
    for p in projects:
        if not p.get("id"):
            raise ValueError(f"Spec project missing 'id': {p}")
        if not p.get("inputs"):
            raise ValueError(f"Spec project {p['id']!r} has no 'inputs'")

        for i, inp in enumerate(p["inputs"]):
            if not isinstance(inp, str):
                raise ValueError(
                    f"Project {p['id']!r} input[{i}] must be a string "
                    f"(folder path or http(s):// URL), got {type(inp).__name__}")

        has_folder = any(not _is_url(i) for i in p["inputs"])
        if not has_folder and not p.get("name"):
            raise ValueError(f"Project {p['id']!r} has no local folder → 'name' is required")
    return projects


def main() -> None:
    ap = argparse.ArgumentParser(description="Extract text from project folders and/or URLs.")
    ap.add_argument("folders", nargs="*", type=Path,
                    help="One or more project folders (backward compat, no URLs).")
    ap.add_argument("--spec", type=Path, default=None,
                    help="JSON spec file with mixed folder+URL inputs (see docstring).")
    ap.add_argument("--output", "-o", type=Path, default=None,
                    help="Write JSON to file instead of stdout.")
    args = ap.parse_args()

    if args.spec and args.folders:
        print("[parse_docs] Use EITHER --spec OR positional folders, not both.", file=sys.stderr)
        sys.exit(2)

    if args.spec:
        try:
            specs = _load_spec(args.spec)
        except (json.JSONDecodeError, ValueError) as e:
            print(f"[parse_docs] Spec error: {e}", file=sys.stderr)
            sys.exit(2)
    elif args.folders:
        specs = _spec_from_positional(args.folders)
    else:
        ap.error("Provide either positional folders or --spec.")

    projects = [parse_project(s) for s in specs]
    payload = {"projects": projects}
    text = json.dumps(payload, ensure_ascii=False, indent=2)

    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(text, encoding="utf-8")
        print(f"[parse_docs] wrote {args.output}", file=sys.stderr)
    else:
        sys.stdout.write(text)


if __name__ == "__main__":
    main()
